#!/usr/bin/env python3
"""把 PPTD 页面转换为元素可编辑的 PowerPoint，而不是整页图片。"""

from __future__ import annotations

import argparse
import html
import re
import sys
from pathlib import Path

import yaml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SKILL_SCRIPTS = Path("/Users/mike/.codex/skills/open-kimi-ppt/scripts")
sys.path.insert(0, str(SKILL_SCRIPTS))
from export_pptx import patch_transitions, verify_output  # noqa: E402


SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5
PAGE_WIDTH_PIXELS = 960
PAGE_HEIGHT_PIXELS = 540
SHAPE_MAP = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rightArrow": MSO_SHAPE.RIGHT_ARROW,
    "ellipse": MSO_SHAPE.OVAL,
}


def pixel_to_inches(value: float) -> float:
    return value * SLIDE_WIDTH_INCHES / PAGE_WIDTH_PIXELS


def resolve_color(value: str | None, colors: dict[str, str], default: str = "#000000") -> RGBColor:
    text = value or default
    if text.startswith("$"):
        text = colors[text[1:]]
    text = text.lstrip("#")[:6]
    return RGBColor.from_string(text.upper())


def resolve_text_style(content: dict, theme: dict) -> dict:
    result: dict = {}
    style_reference = content.get("style")
    if isinstance(style_reference, str) and style_reference.startswith("$"):
        result.update(theme.get("textStyles", {}).get(style_reference[1:], {}))
    result.update({key: value for key, value in content.items() if key != "style"})
    return result


def plain_text(value: str) -> str:
    text = value.replace("<br>", "\n").replace("<br/>", "\n")
    text = re.sub(r"</p>\s*<p[^>]*>", "\n", text)
    text = re.sub(r"</?p[^>]*>", "", text)
    text = re.sub(r"<[^>]+>", "", text)
    return html.unescape(text).strip()


def set_text(slide, element: dict, theme: dict) -> None:
    colors = theme.get("colors", {})
    x, y, width, height = element["bounds"]
    box = slide.shapes.add_textbox(
        Inches(pixel_to_inches(x)), Inches(pixel_to_inches(y)),
        Inches(pixel_to_inches(width)), Inches(pixel_to_inches(height)),
    )
    box.name = element["elementId"]
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = element.get("content", {}).get("wrap", True)
    frame.margin_left = frame.margin_right = Pt(2)
    frame.margin_top = frame.margin_bottom = Pt(1)
    style = resolve_text_style(element["content"], theme)
    align = style.get("align", ["left", "top"])
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM,
    }.get(align[1], MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
    }.get(align[0], PP_ALIGN.LEFT)
    paragraph.line_spacing = style.get("lineHeight", 1.0)
    run = paragraph.add_run()
    run.text = plain_text(str(style.get("text", "")))
    run.font.name = style.get("fontFamily", "Arial")
    run.font.size = Pt(float(style.get("fontSize", 18)) * 0.75)
    run.font.bold = bool(style.get("bold", False))
    run.font.italic = bool(style.get("italic", False))
    run.font.color.rgb = resolve_color(style.get("color"), colors)


def set_shape(slide, element: dict, theme: dict) -> None:
    colors = theme.get("colors", {})
    x, y, width, height = element["bounds"]
    shape = slide.shapes.add_shape(
        SHAPE_MAP.get(element.get("shapeName"), MSO_SHAPE.RECTANGLE),
        Inches(pixel_to_inches(x)), Inches(pixel_to_inches(y)),
        Inches(pixel_to_inches(width)), Inches(pixel_to_inches(height)),
    )
    shape.name = element["elementId"]
    fill = element.get("fill")
    if fill and fill.get("type") == "solid":
        shape.fill.solid()
        shape.fill.fore_color.rgb = resolve_color(fill.get("color"), colors)
    else:
        shape.fill.background()
    border = element.get("border")
    if border:
        shape.line.color.rgb = resolve_color(border.get("color"), colors)
        shape.line.width = Pt(float(border.get("width", 1)) * 0.75)
    else:
        shape.line.fill.background()


def add_picture_fitted(slide, image_path: Path, bounds: list[float], mode: str, name: str) -> None:
    x, y, width, height = bounds
    left, top = Inches(pixel_to_inches(x)), Inches(pixel_to_inches(y))
    target_width, target_height = Inches(pixel_to_inches(width)), Inches(pixel_to_inches(height))
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    if mode == "fill":
        picture = slide.shapes.add_picture(str(image_path), left, top, target_width, target_height)
    elif mode == "contain":
        scale = min(target_width / image_width, target_height / image_height)
        actual_width, actual_height = int(image_width * scale), int(image_height * scale)
        picture = slide.shapes.add_picture(
            str(image_path), left + (target_width - actual_width) // 2,
            top + (target_height - actual_height) // 2, actual_width, actual_height,
        )
    else:
        picture = slide.shapes.add_picture(str(image_path), left, top, target_width, target_height)
        source_ratio = image_width / image_height
        target_ratio = target_width / target_height
        if source_ratio > target_ratio:
            visible = target_ratio / source_ratio
            picture.crop_left = picture.crop_right = (1 - visible) / 2
        else:
            visible = source_ratio / target_ratio
            picture.crop_top = picture.crop_bottom = (1 - visible) / 2
    picture.name = name


def add_background(slide, background: dict, root: Path, theme: dict) -> None:
    colors = theme.get("colors", {})
    if background.get("type") == "image":
        add_picture_fitted(slide, root / background["src"], [0, 0, 960, 540], "fill", "template-background")
    else:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = resolve_color(background.get("color"), colors, "#FFFFFF")


def add_notes(slide, notes: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.clear()
    frame.paragraphs[0].text = notes.strip()


def build(project_directory: Path, output_path: Path) -> dict:
    manifest_path = project_directory / "entering-las-meninas-10-minute-presentation.pptd"
    manifest = yaml.safe_load(manifest_path.read_text(encoding="utf-8"))
    theme = manifest.get("theme", {})
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
    presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    blank = presentation.slide_layouts[6]
    presentation.core_properties.title = manifest["title"]
    presentation.core_properties.author = "Zhuang Chengbo"

    for page_reference in manifest["pages"]:
        page = yaml.safe_load((project_directory / page_reference).read_text(encoding="utf-8"))
        slide = presentation.slides.add_slide(blank)
        add_background(slide, page.get("background", {"type": "solid", "color": "#FFFFFF"}), project_directory, theme)
        for element in page["elements"]:
            # 演讲提示只属于备注，绝不进入观众可见页面。
            if element["elementId"].startswith("cue-"):
                continue
            element_type = element["elementType"]
            if element_type == "text":
                set_text(slide, element, theme)
            elif element_type == "shape":
                set_shape(slide, element, theme)
            elif element_type == "image":
                add_picture_fitted(
                    slide, project_directory / element["src"], element["bounds"],
                    element.get("fit", {}).get("mode", "cover"), element["elementId"],
                )
        add_notes(slide, page["notes"])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    patch_transitions(output_path, "fade")
    return verify_output(output_path, "fade", expect_fonts=False)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("project_directory", type=Path)
    parser.add_argument("output", type=Path)
    arguments = parser.parse_args()
    print(build(arguments.project_directory.resolve(), arguments.output.resolve()))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
